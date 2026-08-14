"""Prueba del motor de plantillas: crea una plantilla .pptx mínima con
banderas de ejemplo, y la rellena usando tus datos reales de output/.
No necesitas tu plantilla real todavía — esto valida que el mecanismo
de reemplazo funciona antes de conectar tu archivo de verdad.

Uso: python scripts/test_ppt_core.py 2026-08-12 BSNC
"""
import sys
from pathlib import Path

sys.path.append(str(Path(__file__).resolve().parent.parent))

from pptx import Presentation
from pptx.util import Inches, Pt

from config.settings import OUTPUT_DIR
from models.metrica import MetricaRepository
from models.entidad import EntidadRepository
from services.ppt.ppt_core import generar_ppt_comparativo
from services.ppt.datos_ppt import buscar_archivo_anterior


def crear_plantilla_ejemplo(ruta: Path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    textos = [
        "Fecha: <<fecha>>",
        "Push Security: <<push_security>>",
        "Cumplimiento: <<push_security_cumplimiento>>",
        "Umbral: <<push_security_umbral>>  Num: <<push_security_numerador>>  Den: <<push_security_denominador>>",
        "Cifrado: <<cifrado>>  |  Status: <<cifrado_status>>",
        "IOS desactualizados: <<moviles_compliant(IOS)>>",
        "Android desactualizados: <<moviles_compliant(Android)>>",
    ]
    y = Inches(0.5)
    for texto in textos:
        box = slide.shapes.add_textbox(Inches(0.5), y, Inches(9), Inches(0.5))
        box.text_frame.text = texto
        for p in box.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(14)
        y += Inches(0.6)

    prs.save(ruta)


def main():
    fecha = sys.argv[1] if len(sys.argv) > 1 else None
    entidad_nombre = sys.argv[2] if len(sys.argv) > 2 else "BSNC"
    if fecha is None:
        print("❌ Uso: python scripts/test_ppt_core.py <fecha> <entidad>")
        return

    repo_entidad = EntidadRepository()
    entidad = repo_entidad.obtener_por_nombre(entidad_nombre)
    if entidad is None:
        print(f"❌ No existe la entidad '{entidad_nombre}'")
        return

    nombre_archivo = entidad.nombre_archivo_salida or entidad.nombre
    ruta_excel_actual = OUTPUT_DIR / fecha / f"{nombre_archivo}.xlsx"
    if not ruta_excel_actual.exists():
        print(f"❌ No existe {ruta_excel_actual}")
        return

    ruta_excel_anterior = buscar_archivo_anterior(nombre_archivo, fecha)
    ruta_cumplimiento = OUTPUT_DIR / fecha / f"resumen_cumplimiento_{entidad.nombre.lower()}.xlsx"

    ruta_plantilla = Path("plantilla_ejemplo_temp.pptx")
    crear_plantilla_ejemplo(ruta_plantilla)
    print(f"📄 Plantilla de ejemplo creada: {ruta_plantilla}")

    metricas = MetricaRepository().listar()
    ruta_salida = Path("output_test_ppt_core.pptx")

    n = generar_ppt_comparativo(
        ruta_excel_actual, ruta_excel_anterior, ruta_plantilla, ruta_salida,
        metricas, ruta_resumen_cumplimiento=ruta_cumplimiento)

    print(f"\n✅ Generado: {ruta_salida} ({n} banderas reemplazadas)")
    print("   Ábrelo y compáralo con la plantilla de ejemplo para confirmar que los valores se ven bien.")


if __name__ == "__main__":
    main()