"""Reemplaza banderas <<...>> dentro de un text_frame de PowerPoint,
manteniendo el formato/color especial según el tipo de bandera (normal,
cumplimiento, umbral/numerador/denominador, u operación con paréntesis
como <<hoja(operacion)>>). Puerto de ppt_bandera_handler.py del sistema
original, adaptado para usar Metrica.operaciones_ppt en vez de las
'instrucciones' derivadas de reportes.json.
"""
import re
import pandas as pd
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from models.metrica import Metrica
from services.reglas import evaluar_grupos
from services.ppt.plantilla_utils import formato_cambio
from services.ppt.banderas import normalizar_clave


# ══════════════════════ ÍNDICE DE OPERACIONES ══════════════════════

def construir_indice_operaciones(metricas: list[Metrica]) -> dict:
    """{'hoja_normalizada': [{'nombre': str, 'condicion': [[Condicion]]}]}
    agregando las operaciones_ppt de TODAS las métricas, agrupadas por
    la hoja de la que leen (hoja_fuente). También registra cada
    operación bajo su alias legacy (ej. 'qualys' además de
    'qualys_no_cmdb'), para compatibilidad con la plantilla original."""
    from services.ppt.legacy_alias import ALIAS_OPERACION_HOJA

    alias_por_hoja_real: dict[str, list[str]] = {}
    for alias, nombre_real in ALIAS_OPERACION_HOJA.items():
        alias_por_hoja_real.setdefault(normalizar_clave(nombre_real), []).append(alias)

    indice: dict[str, list] = {}
    for m in metricas:
        for op in m.operaciones_ppt:
            clave_hoja = normalizar_clave(op.hoja_fuente)
            entrada = {"nombre": op.nombre, "condicion": op.condicion}
            indice.setdefault(clave_hoja, []).append(entrada)
            for alias in alias_por_hoja_real.get(clave_hoja, []):
                indice.setdefault(alias, []).append(entrada)
    return indice


def _leer_hoja_normalizada(ruta_excel, hoja_normalizada: str) -> pd.DataFrame | None:
    try:
        xl = pd.ExcelFile(ruta_excel)
    except Exception:
        return None
    for hoja in xl.sheet_names:
        if normalizar_clave(hoja) == hoja_normalizada:
            try:
                return xl.parse(hoja)
            except Exception:
                return None
    return None


def aplicar_operacion(ruta_excel_actual, hoja_normalizada: str, operacion_nombre: str,
                      operaciones_index: dict, ruta_excel_anterior=None) -> str:
    """Cuenta filas que cumplen la condición de 'operacion_nombre' en la
    hoja correspondiente, semana actual vs anterior, y retorna el texto
    formateado (mismo formato que una bandera normal)."""
    from services.ppt.legacy_alias import ALIAS_OPERACION_HOJA

    ops = operaciones_index.get(hoja_normalizada, [])
    op = next((o for o in ops if o["nombre"].lower() == operacion_nombre.lower()), None)
    if op is None:
        return f"⚠️ Operación '{operacion_nombre}' no configurada"

    # 'hoja_normalizada' puede ser un alias corto (ej. 'qualys') — para
    # LEER los datos hace falta el nombre real de la hoja en el Excel.
    nombre_real = ALIAS_OPERACION_HOJA.get(hoja_normalizada)
    hoja_para_leer = normalizar_clave(nombre_real) if nombre_real else hoja_normalizada

    df_actual = _leer_hoja_normalizada(ruta_excel_actual, hoja_para_leer)
    if df_actual is None:
        return f"⚠️ Hoja '{hoja_normalizada}' no encontrada"

    conteo_actual = int(evaluar_grupos(df_actual, op["condicion"]).sum())

    conteo_anterior = None
    if ruta_excel_anterior:
        df_anterior = _leer_hoja_normalizada(ruta_excel_anterior, hoja_para_leer)
        if df_anterior is not None:
            conteo_anterior = int(evaluar_grupos(df_anterior, op["condicion"]).sum())

    if conteo_anterior is not None:
        texto, _color, _ = formato_cambio(conteo_actual, conteo_anterior)
        return texto
    return str(conteo_actual)


# ══════════════════════ REEMPLAZO CON FORMATO ══════════════════════

def reemplazar_con_formato(text_frame, texto_original: str, valores: dict,
                           ruta_excel_actual=None, ruta_excel_anterior=None,
                           operaciones_index: dict = None):
    """Reemplaza todas las banderas de un text_frame, reconstruyendo el
    párrafo con el color/formato correspondiente a cada tipo de bandera."""
    operaciones_index = operaciones_index or {}
    valores_norm = {re.sub(r"[\s\-\.]+", "_", str(k).lower()): v for k, v in valores.items()}

    banderas = []
    for match in re.finditer(r"<<([^>]+)>>", texto_original):
        bandera_texto = match.group(1)
        op_match = re.match(r"^([^\(]+)\(([^\)]+)\)$", bandera_texto)

        if op_match:
            nombre_hoja = normalizar_clave(op_match.group(1).strip())
            nombre_operacion = op_match.group(2).strip()
            valor_operacion = aplicar_operacion(
                ruta_excel_actual, nombre_hoja, nombre_operacion,
                operaciones_index, ruta_excel_anterior=ruta_excel_anterior)
            banderas.append({
                "start": match.start(), "end": match.end(),
                "bandera": bandera_texto, "bandera_norm": nombre_hoja,
                "valor": valor_operacion, "es_operacion": True,
            })
        else:
            bandera_norm = normalizar_clave(bandera_texto)
            if bandera_norm in valores_norm:
                banderas.append({
                    "start": match.start(), "end": match.end(),
                    "bandera": bandera_texto, "bandera_norm": bandera_norm,
                    "valor": str(valores_norm[bandera_norm]), "es_operacion": False,
                })

    if not banderas:
        return 0

    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER

    pos_actual = 0
    for band_info in banderas:
        if pos_actual < band_info["start"]:
            run_antes = p.add_run()
            run_antes.text = texto_original[pos_actual:band_info["start"]]
            run_antes.font.size = Pt(10)

        valor = band_info["valor"]

        if band_info["es_operacion"] or _es_comparacion_simple(valor):
            _escribir_comparacion(p, valor)
        elif band_info["bandera_norm"].endswith("_cumplimiento"):
            _escribir_cumplimiento(p, valor)
        elif band_info["bandera_norm"].endswith(("_umbral", "_denominador", "_numerador")):
            run_valor = p.add_run()
            run_valor.text = str(valor)
            run_valor.font.size = Pt(10)
            run_valor.font.color.rgb = RGBColor(0, 0, 0)
        elif band_info["bandera_norm"].endswith("_status"):
            run_valor = p.add_run()
            run_valor.text = str(valor)
            run_valor.font.bold = True
            run_valor.font.size = Pt(10)
            run_valor.font.color.rgb = _determinar_color_status(valor)
        else:
            run_valor = p.add_run()
            run_valor.text = str(valor)
            run_valor.font.size = Pt(10)

        pos_actual = band_info["end"]

    if pos_actual < len(texto_original):
        run_restante = p.add_run()
        run_restante.text = texto_original[pos_actual:]
        run_restante.font.size = Pt(10)

    return len(banderas)


def _es_comparacion_simple(valor: str) -> bool:
    return bool(re.match(r"^\s*[=↑↓]\s*\d+", str(valor)))


def _escribir_comparacion(p, valor: str):
    m = re.match(r"^\s*([=↑↓]\s*\d+)\s*(.*)$", valor)
    if m:
        parte_color, parte_texto = m.group(1), m.group(2)
    else:
        parte_color, parte_texto = valor, ""

    simbolo = parte_color[0]
    numero = parte_color[1:]

    run_simbolo = p.add_run()
    run_simbolo.text = simbolo
    run_simbolo.font.bold = True
    run_simbolo.font.size = Pt(10)
    color_simbolo = _determinar_color_simbolo(simbolo)
    if color_simbolo:
        run_simbolo.font.color.rgb = color_simbolo

    run_numero = p.add_run()
    run_numero.text = numero
    run_numero.font.bold = True
    run_numero.font.size = Pt(10)
    if color_simbolo:
        run_numero.font.color.rgb = color_simbolo

    if parte_texto:
        run_desc = p.add_run()
        run_desc.text = " " + parte_texto
        run_desc.font.size = Pt(9)


def _escribir_cumplimiento(p, valor: str):
    partes = [x.strip() for x in str(valor).split("/", 1)]
    cumplimiento_txt = partes[0] if partes and partes[0] else "N/A"
    status_txt = partes[1] if len(partes) > 1 and partes[1] else "N/A"

    run_c = p.add_run()
    run_c.text = cumplimiento_txt
    run_c.font.size = Pt(10)

    run_sep = p.add_run()
    run_sep.text = " / "
    run_sep.font.size = Pt(10)

    run_s = p.add_run()
    run_s.text = status_txt
    run_s.font.bold = True
    run_s.font.size = Pt(10)
    run_s.font.color.rgb = _determinar_color_status(status_txt)


def _determinar_color_status(status: str):
    status_norm = str(status).strip().lower()
    if status_norm in ("ok", "cumple"):
        return RGBColor(0, 128, 0)
    if status_norm in ("failed", "no cumple", "nocumple"):
        return RGBColor(192, 0, 0)
    return RGBColor(0, 0, 0)


def _determinar_color_simbolo(simbolo: str):
    if simbolo == "↑":
        return RGBColor(192, 0, 0)
    if simbolo == "↓":
        return RGBColor(0, 128, 0)
    if simbolo == "=":
        return RGBColor(0, 112, 192)
    return None


def reemplazar_simple(text: str, valores: dict) -> str:
    """Reemplazo de banderas sin formato de color (fallback si algo falla
    con reemplazar_con_formato)."""
    valores_norm = {re.sub(r"[\s\-\.]+", "_", str(k).lower()): v for k, v in valores.items()}

    def repl(m):
        bandera_norm = re.sub(r"[\s\-\.]+", "_", m.group(1).lower())
        return str(valores_norm.get(bandera_norm, m.group(0)))

    return re.sub(r"<<(.*?)>>", repl, text)