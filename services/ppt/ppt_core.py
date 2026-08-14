"""Motor principal de generación de PPT por plantilla: abre la
plantilla, recorre cada diapositiva/shape buscando banderas <<...>>, y
las reemplaza con los valores calculados. Puerto de ppt_core.py del
sistema original — aquí NO hace falta el diccionario de mapeo manual
(MAPEO_RESUMEN_CUMPLIMIENTO) porque el nombre de la hoja del reporte y
el nombre en cumplimiento son siempre el mismo: Metrica.nombre.
"""
from pathlib import Path
from datetime import datetime

from pptx import Presentation

from models.metrica import Metrica
from services.ppt.plantilla_utils import leer_hojas_excel, comparar_excel, buscar_bandera
from services.ppt.bandera_handler import (reemplazar_con_formato, reemplazar_simple,
                                          construir_indice_operaciones)
from services.ppt.banderas import normalizar_clave
from services.ppt.datos_ppt import cargar_cumplimiento
from services.ppt.legacy_alias import ALIAS_HOJA, ALIAS_CUMPLIMIENTO


def _construir_valores(metricas: list[Metrica], comparativo: dict,
                       cumplimiento_map: dict, fecha_texto: str, titulo: str) -> dict:
    """Arma el diccionario {bandera: valor} con TODAS las banderas
    posibles: comparación semanal (por hoja) + cumplimiento (por métrica)."""
    valores = {"fecha": fecha_texto, "titulo": titulo}

    for hoja, datos in comparativo.items():
        valores[normalizar_clave(hoja)] = datos["texto"]

    for m in metricas:
        if not m.cumplimiento.aplica:
            continue
        clave = normalizar_clave(m.nombre)
        info = cumplimiento_map.get(m.nombre.strip().lower())
        if info is None:
            continue
        resultado = info.get("resultado")
        resultado_pct = f"{resultado * 100:.1f}%" if isinstance(resultado, (int, float)) else "N/A"
        status = info.get("estado", "N/A")
        valores[f"{clave}_cumplimiento"] = f"{resultado_pct} / {status}"
        valores[f"{clave}_status"] = status
        umbral = info.get("umbral")
        valores[f"{clave}_umbral"] = f"{umbral * 100:.0f}%" if isinstance(umbral, (int, float)) else "N/A"
        valores[f"{clave}_numerador"] = str(info.get("numerador", "N/A"))
        valores[f"{clave}_denominador"] = str(info.get("denominador", "N/A"))

    # ── Alias de compatibilidad con la plantilla original ──
    for alias, nombre_real in ALIAS_HOJA.items():
        clave_real = normalizar_clave(nombre_real)
        if clave_real in valores:
            valores[alias] = valores[clave_real]

    for alias, nombre_real in ALIAS_CUMPLIMIENTO.items():
        clave_real = normalizar_clave(nombre_real)
        for sufijo in ("_cumplimiento", "_status", "_umbral", "_numerador", "_denominador"):
            k = f"{clave_real}{sufijo}"
            if k in valores:
                valores[f"{alias}{sufijo}"] = valores[k]

    return valores


def generar_ppt_comparativo(ruta_excel_actual, ruta_excel_anterior, ruta_plantilla,
                            ruta_salida, metricas: list[Metrica],
                            ruta_resumen_cumplimiento=None) -> int:
    """
    Abre ruta_plantilla, reemplaza todas las banderas encontradas, y
    guarda el resultado en ruta_salida. Retorna cuántas banderas se
    reemplazaron en total.
    """
    print("\n📊 GENERANDO PRESENTACIÓN COMPARATIVA")
    resumenes_actual = leer_hojas_excel(str(ruta_excel_actual))
    resumenes_anterior = leer_hojas_excel(str(ruta_excel_anterior)) if ruta_excel_anterior else {}
    print(f"   Hojas: {len(resumenes_actual)} actual / {len(resumenes_anterior)} anterior")

    comparativo = comparar_excel(resumenes_actual, resumenes_anterior)
    cumplimiento_map = cargar_cumplimiento(Path(ruta_resumen_cumplimiento)) if ruta_resumen_cumplimiento else {}
    operaciones_index = construir_indice_operaciones(metricas)

    fecha_texto = datetime.now().strftime("%d/%m/%Y")
    titulo = Path(ruta_excel_actual).stem
    valores = _construir_valores(metricas, comparativo, cumplimiento_map, fecha_texto, titulo)

    prs = Presentation(str(ruta_plantilla))
    print(f"   {len(prs.slides)} diapositiva(s) en plantilla")

    banderas_reemplazadas = 0
    for idx_slide, slide in enumerate(prs.slides, 1):
        banderas_en_slide = 0
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            texto_original = shape.text
            if not buscar_bandera(texto_original):
                continue
            if hasattr(shape, "text_frame"):
                try:
                    n = reemplazar_con_formato(
                        shape.text_frame, texto_original, valores,
                        ruta_excel_actual=str(ruta_excel_actual),
                        ruta_excel_anterior=str(ruta_excel_anterior) if ruta_excel_anterior else None,
                        operaciones_index=operaciones_index)
                    banderas_en_slide += n
                except Exception as e:
                    print(f"      ⚠️  Fallback a reemplazo simple ({e})")
                    nuevo = reemplazar_simple(texto_original, valores)
                    if nuevo != texto_original:
                        shape.text = nuevo
                        banderas_en_slide += len(buscar_bandera(texto_original))
            else:
                nuevo = reemplazar_simple(texto_original, valores)
                if nuevo != texto_original:
                    shape.text = nuevo
                    banderas_en_slide += len(buscar_bandera(texto_original))

        if banderas_en_slide:
            print(f"   [{idx_slide}] {banderas_en_slide} bandera(s) reemplazada(s)")
        banderas_reemplazadas += banderas_en_slide

    prs.save(str(ruta_salida))
    print(f"\n✅ PPT guardado: {ruta_salida} ({banderas_reemplazadas} banderas reemplazadas)")
    return banderas_reemplazadas