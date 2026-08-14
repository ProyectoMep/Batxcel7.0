"""Diseño visual de la presentación PowerPoint. Este es el ÚNICO archivo
que hay que tocar si se quiere cambiar cómo se ve el PPT (colores,
tamaños, tipografía, layout de las tarjetas). No contiene lógica de
negocio — solo recibe datos ya calculados y los dibuja.

Las tarjetas usan bloques de texto (etiqueta pequeña + valor grande),
no tablas nativas de PowerPoint, para tener control total del espaciado
y evitar que se vean apretadas.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ══════════════════════ PALETA Y TIPOGRAFÍA ══════════════════════

ANCHO_SLIDE = Inches(13.333)
ALTO_SLIDE = Inches(7.5)

FUENTE = "Calibri"

ROJO_SANTANDER = RGBColor(0xEC, 0x00, 0x00)
AZUL_TEXTO = RGBColor(0x1F, 0x3B, 0x57)
NAVY_BARRA = RGBColor(0x1B, 0x2A, 0x4A)
GRIS_CLARO = RGBColor(0xF5, 0xF6, 0xF8)
GRIS_BORDE = RGBColor(0xDD, 0xE1, 0xE6)
GRIS_ETIQUETA = RGBColor(0x8A, 0x93, 0x9E)
BLANCO = RGBColor(0xFF, 0xFF, 0xFF)
NEGRO = RGBColor(0x1A, 0x1A, 0x1A)

VERDE_CUMPLE = RGBColor(0x1E, 0x8E, 0x3E)
VERDE_FONDO = RGBColor(0xE8, 0xF5, 0xEA)
ROJO_NO_CUMPLE = RGBColor(0xC5, 0x22, 0x1F)
ROJO_FONDO = RGBColor(0xFC, 0xEA, 0xEA)

ROJO_AUMENTO = RGBColor(0xC5, 0x22, 0x1F)        # más incumplimientos = malo
VERDE_DISMINUCION = RGBColor(0x1E, 0x8E, 0x3E)   # menos incumplimientos = bueno
AZUL_SIN_CAMBIO = RGBColor(0x1A, 0x73, 0xC7)


# ══════════════════════ UTILIDADES BASE ══════════════════════

def _set_text(shape, texto, tamano=11, negrita=False, color=NEGRO,
              alineacion=PP_ALIGN.LEFT, fuente=FUENTE, anchor=None):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(1)
    tf.margin_right = Pt(1)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    if anchor is not None:
        tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = alineacion
    run = p.add_run()
    run.text = texto
    run.font.size = Pt(tamano)
    run.font.bold = negrita
    run.font.color.rgb = color
    run.font.name = fuente
    return shape


def _rect(slide, x, y, w, h, color_relleno=None, color_borde=None, redondeado=False, grosor_borde=0.75):
    forma = MSO_SHAPE.ROUNDED_RECTANGLE if redondeado else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(forma, x, y, w, h)
    shape.shadow.inherit = False
    if redondeado:
        try:
            shape.adjustments[0] = 0.06
        except Exception:
            pass
    if color_relleno is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = color_relleno
    else:
        shape.fill.background()
    if color_borde is not None:
        shape.line.color.rgb = color_borde
        shape.line.width = Pt(grosor_borde)
    else:
        shape.line.fill.background()
    return shape


def _caja_texto(slide, x, y, w, h, texto, **kwargs):
    box = slide.shapes.add_textbox(x, y, w, h)
    _set_text(box, texto, **kwargs)
    return box


def _bloque_estadistica(slide, x, y, w, etiqueta, valor, color_valor=NEGRO,
                        tamano_etiqueta=7, tamano_valor=12.5):
    """Un mini bloque 'KPI': etiqueta pequeña gris arriba, valor grande
    y en negrita abajo, ambos centrados."""
    _caja_texto(slide, x, y, w, Inches(0.14), etiqueta.upper(), tamano=tamano_etiqueta,
               negrita=True, color=GRIS_ETIQUETA, alineacion=PP_ALIGN.CENTER)
    _caja_texto(slide, x, y + Inches(0.15), w, Inches(0.26), valor, tamano=tamano_valor,
               negrita=True, color=color_valor, alineacion=PP_ALIGN.CENTER)


def _badge_resultado(slide, x, y, w, h, resultado_pct: str, estado: str):
    """Insignia de resultado: fondo tintado (verde/rojo suave), % grande
    y estado debajo, todo centrado verticalmente."""
    es_cumple = estado == "Cumple"
    color_texto = VERDE_CUMPLE if es_cumple else ROJO_NO_CUMPLE
    color_fondo = VERDE_FONDO if es_cumple else ROJO_FONDO

    _rect(slide, x, y, w, h, color_relleno=color_fondo, redondeado=True)
    caja = slide.shapes.add_textbox(x, y, w, h)
    tf = caja.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = resultado_pct
    r1.font.size = Pt(15)
    r1.font.bold = True
    r1.font.name = FUENTE
    r1.font.color.rgb = color_texto
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = estado
    r2.font.size = Pt(9.5)
    r2.font.bold = True
    r2.font.name = FUENTE
    r2.font.color.rgb = color_texto


# ══════════════════════ REGLA DE COLOR PARA COMPARACIONES ══════════════════════

def formatear_comparacion(actual: int, anterior: int) -> tuple[str, "RGBColor"]:
    """Rojo = aumentaron los incumplimientos (malo). Verde = disminuyeron
    (bueno). Azul = sin cambio."""
    diferencia = actual - anterior
    if diferencia > 0:
        return f"↑ {actual} vs. semana pasada ({anterior})", ROJO_AUMENTO
    if diferencia < 0:
        return f"↓ {actual} vs. semana pasada ({anterior})", VERDE_DISMINUCION
    return f"= {actual} sin cambios vs. semana pasada", AZUL_SIN_CAMBIO


# ══════════════════════ PRESENTACIÓN BASE ══════════════════════

def crear_presentacion() -> Presentation:
    prs = Presentation()
    prs.slide_width = ANCHO_SLIDE
    prs.slide_height = ALTO_SLIDE
    return prs


def nueva_slide_en_blanco(prs):
    layout_blanco = prs.slide_layouts[6]
    return prs.slides.add_slide(layout_blanco)


# ══════════════════════ PORTADA ══════════════════════

def agregar_portada(prs, entidad_nombre: str, fecha_texto: str, descripcion: str,
                    ruta_logo_color: str = None):
    slide = nueva_slide_en_blanco(prs)

    _caja_texto(slide, Inches(4.2), Inches(0.3), Inches(8.5), Inches(0.4),
               f"Gerencia de Ciberseguridad     Bogotá     {fecha_texto}",
               tamano=11, color=AZUL_TEXTO, alineacion=PP_ALIGN.RIGHT)

    if ruta_logo_color:
        try:
            slide.shapes.add_picture(ruta_logo_color, Inches(0.6), Inches(0.9), height=Inches(0.55))
        except Exception:
            pass

    _caja_texto(slide, Inches(0.6), Inches(2.6), Inches(6), Inches(0.4),
               f"Resumen Cyber {entidad_nombre}",
               tamano=18, negrita=True, color=ROJO_SANTANDER)

    _caja_texto(slide, Inches(0.6), Inches(3.4), Inches(5.8), Inches(1.6),
               descripcion, tamano=18, color=NEGRO)

    return slide


# ══════════════════════ BARRA DE SECCIÓN ══════════════════════

ICONOS_SECCION = {
    "Agentes - Workstations": "🖧",
    "Herramientas - Workstations": "⚙",
    "Usuarios": "👥",
    "Antimalware - Dispositivos Moviles": "🔌",
}


def agregar_barra_seccion(slide, texto: str, y, x=Inches(0.5), w=Inches(12.33), h=Inches(0.4)):
    _rect(slide, x, y, w, h, color_relleno=NAVY_BARRA, redondeado=True)
    icono = ICONOS_SECCION.get(texto, "")
    texto_final = f"{icono}   {texto}" if icono else texto
    _caja_texto(slide, x + Inches(0.2), y, w - Inches(0.4), h, texto_final,
               tamano=13.5, negrita=True, color=BLANCO, anchor=MSO_ANCHOR.MIDDLE)
    return y + h + Inches(0.16)


def agregar_encabezado_slide(slide, titulo_reporte: str, fecha_texto: str):
    _caja_texto(slide, Inches(0.5), Inches(0.18), Inches(6), Inches(0.3),
               titulo_reporte, tamano=11, negrita=True, color=NEGRO)
    _caja_texto(slide, Inches(10.4), Inches(0.18), Inches(2.4), Inches(0.3),
               fecha_texto, tamano=10, negrita=True, color=ROJO_SANTANDER,
               alineacion=PP_ALIGN.RIGHT)


# ══════════════════════ TARJETA DE MÉTRICA ══════════════════════

ANCHO_BADGE = Inches(1.35)


def dibujar_tarjeta(slide, x, y, w, h, titulo: str, dato_cumplimiento: dict):
    """
    dato_cumplimiento: {
      'comparacion_texto': str, 'comparacion_color': RGBColor,
      'umbral': float, 'numerador': int, 'denominador': int,
      'resultado_pct': str, 'estado': 'Cumple'|'No Cumple'|'Sin datos'
    }
    Dibuja: fondo blanco con acento de color -> título -> comparación ->
    3 bloques KPI (Umbral/Numerador/Denominador) -> insignia de resultado.
    """
    color_estado = (VERDE_CUMPLE if dato_cumplimiento["estado"] == "Cumple"
                    else ROJO_NO_CUMPLE if dato_cumplimiento["estado"] == "No Cumple"
                    else GRIS_ETIQUETA)

    _rect(slide, x, y, w, h, color_relleno=BLANCO, color_borde=GRIS_BORDE, redondeado=True)
    _rect(slide, x, y + Inches(0.06), Inches(0.06), h - Inches(0.12), color_relleno=color_estado, redondeado=True)

    pad = Inches(0.18)
    x_c = x + pad + Inches(0.05)
    w_c = w - pad - Inches(0.05) - ANCHO_BADGE - Inches(0.15)

    _caja_texto(slide, x_c, y + Inches(0.1), w_c, Inches(0.22), titulo,
               tamano=11.5, negrita=True, color=AZUL_TEXTO)
    _caja_texto(slide, x_c, y + Inches(0.32), w_c, Inches(0.2),
               dato_cumplimiento["comparacion_texto"], tamano=8.5, negrita=True,
               color=dato_cumplimiento["comparacion_color"])

    y_kpi = y + h - Inches(0.46)
    ancho_kpi = w_c / 3
    _bloque_estadistica(slide, x_c, y_kpi, ancho_kpi, "Umbral", f"{dato_cumplimiento['umbral'] * 100:.0f}%")
    _bloque_estadistica(slide, x_c + ancho_kpi, y_kpi, ancho_kpi, "Numer.", str(dato_cumplimiento["numerador"]))
    _bloque_estadistica(slide, x_c + ancho_kpi * 2, y_kpi, ancho_kpi, "Denom.", str(dato_cumplimiento["denominador"]))

    x_badge = x + w - ANCHO_BADGE - Inches(0.12)
    _badge_resultado(slide, x_badge, y + Inches(0.1), ANCHO_BADGE, h - Inches(0.2),
                     dato_cumplimiento["resultado_pct"], dato_cumplimiento["estado"])


# ══════════════════════ GRUPO VERTICAL (etiqueta compartida) ══════════════════════

ANCHO_ETIQUETA_VERTICAL = Inches(0.32)


def dibujar_grupo_vertical(slide, x, y, w, h, etiqueta: str, subtarjetas: list[dict]):
    """Dibuja una barra vertical navy con el texto rotado, y a la derecha
    apila las subtarjetas (todas del mismo alto)."""
    _rect(slide, x, y, ANCHO_ETIQUETA_VERTICAL, h, color_relleno=NAVY_BARRA, redondeado=True)
    caja_etiqueta = slide.shapes.add_textbox(x - (h / 2) + (ANCHO_ETIQUETA_VERTICAL / 2),
                                             y + (h / 2) - Inches(0.15),
                                             h, ANCHO_ETIQUETA_VERTICAL)
    caja_etiqueta.rotation = 270
    tf = caja_etiqueta.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = etiqueta
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.name = FUENTE
    run.font.color.rgb = BLANCO

    x_tarjetas = x + ANCHO_ETIQUETA_VERTICAL + Inches(0.1)
    w_tarjetas = w - ANCHO_ETIQUETA_VERTICAL - Inches(0.1)
    alto_cada_una = h / len(subtarjetas)

    for i, sub in enumerate(subtarjetas):
        y_sub = y + (alto_cada_una * i)
        dibujar_tarjeta(slide, x_tarjetas, y_sub, w_tarjetas, alto_cada_una,
                        sub["titulo"], sub["dato"])


# ══════════════════════ CAJA DESTACADA (métricas "sueltas") ══════════════════════

def dibujar_divisor_vertical(slide, x, y, h):
    linea = slide.shapes.add_connector(1, x, y, x, y + h)
    linea.line.color.rgb = GRIS_BORDE
    linea.line.width = Pt(1)


def dibujar_caja_destacada(slide, x, y, w, h, titulo: str, dato_cumplimiento: dict,
                           resultado_a_la_izquierda: bool = True):
    """Caja ancha para métricas 'sueltas' (High Risk Software,
    Workstations Retired): título + comparación centrados arriba, 3
    bloques KPI en el medio, insignia de resultado a un lado."""
    ancho_badge = Inches(1.7)

    if resultado_a_la_izquierda:
        x_badge = x + Inches(0.1)
        x_contenido = x + ancho_badge + Inches(0.3)
    else:
        x_contenido = x + Inches(0.1)
        x_badge = x + w - ancho_badge - Inches(0.1)

    w_contenido = w - ancho_badge - Inches(0.4)

    _rect(slide, x, y, w, h, color_relleno=BLANCO, color_borde=GRIS_BORDE, redondeado=True)

    _badge_resultado(slide, x_badge, y + Inches(0.12), ancho_badge, h - Inches(0.24),
                     dato_cumplimiento["resultado_pct"], dato_cumplimiento["estado"])

    _caja_texto(slide, x_contenido, y + Inches(0.14), w_contenido, Inches(0.26), titulo,
               tamano=14.5, negrita=True, color=AZUL_TEXTO, alineacion=PP_ALIGN.CENTER)
    _caja_texto(slide, x_contenido, y + Inches(0.42), w_contenido, Inches(0.22),
               dato_cumplimiento["comparacion_texto"], tamano=9.5, negrita=True,
               color=dato_cumplimiento["comparacion_color"], alineacion=PP_ALIGN.CENTER)

    y_kpi = y + h - Inches(0.5)
    ancho_kpi = w_contenido / 3
    _bloque_estadistica(slide, x_contenido, y_kpi, ancho_kpi, "Umbral",
                        f"{dato_cumplimiento['umbral'] * 100:.0f}%", tamano_valor=14)
    _bloque_estadistica(slide, x_contenido + ancho_kpi, y_kpi, ancho_kpi, "Numerador",
                        str(dato_cumplimiento["numerador"]), tamano_valor=14)
    _bloque_estadistica(slide, x_contenido + ancho_kpi * 2, y_kpi, ancho_kpi, "Denominador",
                        str(dato_cumplimiento["denominador"]), tamano_valor=14)


# ══════════════════════ FILA DE OPERACIÓN (subconteo tipo IOS/Android) ══════════════════════

ALTO_FILA_OPERACION = Inches(0.42)


def dibujar_fila_operacion(slide, x, y, w, h, nombre: str, texto_comparacion: str, color_comparacion):
    """Fila compacta tipo 'pill': nombre a la izquierda sobre fondo
    tenue, comparación a la derecha. Usada para IOS/Android dentro de
    un grupo."""
    _rect(slide, x, y, w, h - Inches(0.06), color_relleno=GRIS_CLARO, redondeado=True)
    mitad = w * 0.38
    _caja_texto(slide, x + Inches(0.12), y, mitad, h - Inches(0.06), nombre,
               tamano=9.5, negrita=True, color=AZUL_TEXTO, anchor=MSO_ANCHOR.MIDDLE)
    _caja_texto(slide, x + mitad, y, w - mitad - Inches(0.12), h - Inches(0.06), texto_comparacion,
               tamano=8.5, negrita=True, color=color_comparacion,
               alineacion=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


# ══════════════════════ CIERRE ══════════════════════

def agregar_cierre(prs, ruta_logo_blanco: str = None):
    slide = nueva_slide_en_blanco(prs)
    _rect(slide, 0, 0, Inches(5.5), ALTO_SLIDE, color_relleno=BLANCO)
    _rect(slide, Inches(5.5), 0, ANCHO_SLIDE - Inches(5.5), ALTO_SLIDE, color_relleno=ROJO_SANTANDER)

    _caja_texto(slide, Inches(0.6), Inches(1.9), Inches(4.5), Inches(1),
               "Gracias.", tamano=44, color=ROJO_SANTANDER)

    _caja_texto(slide, Inches(0.6), Inches(4.0), Inches(4.3), Inches(1.2),
               "Nuestro propósito es ayudar a personas y empresas a prosperar.\n\n"
               "Nuestra cultura se basa en la creencia de que todo lo que hacemos debe ser.",
               tamano=11, color=AZUL_TEXTO)

    _caja_texto(slide, Inches(0.6), Inches(5.5), Inches(4.3), Inches(0.4),
               "Es el momento de avanzar.", tamano=12, negrita=True, color=ROJO_SANTANDER)

    if ruta_logo_blanco:
        try:
            slide.shapes.add_picture(ruta_logo_blanco, Inches(8), Inches(3.3), height=Inches(0.9))
        except Exception:
            pass

    return slide


# ══════════════════════ SLIDE ESTÁTICA: NOMENCLATURA ══════════════════════

TABLA_NOMENCLATURA = [
    ("Empresa", "Valor", "Departamento", "Ejemplo"),
    ("Banco Santander de Colombia", "IOCODC", "1 carácter 0-z", "IOCODCHX268607P"),
    ("SBI - Santander Banca de Inversión", "IOCSBI", "1 carácter 0-z", "IOCSBIHX268607P"),
    ("UNC - Universia Colombia", "IOCUNC", "1 carácter 0-z", "IOCUNCHX268607P"),
    ("CACEIS Colombia", "IOCCDC", "1 carácter 0-z", "IOCCDCHX268607P"),
    ("Santander Consumer Colombia", "SCCOL", "5 caracteres 0-z", "SCCOL00N370358P"),
    ("Santander Financing Colombia", "SFCOL", "5 caracteres 0-z", "SFCOL01N458437P"),
]

NOTA_NOMENCLATURA = ("Nota: Los equipos que están con nomenclatura Santander Consumer, "
                     "deben migrar a Santander Financing.")

TEXTO_INTRO_NOMENCLATURA = (
    "Nomenclatura aplicable a los entornos de Global: SGT, SanHQ, Ucloud, ODS, "
    "PagoNxt, SCEUADFRP, SCGBN, OpenBank, Zurich, SCIB, EUT…"
)

TEXTO_MATRICULA = [
    "7 dígitos alfanuméricos:                    Ejemplo: SA01111N222222P",
    "Debe comenzar por la letra \u201cN\u201d o \u201cX\u201d",
    "Corresponderá siempre a la matrícula de empleado o del usuario "
    "(OBLIGATORIO: 6 DIGITOS A PARTIR DE LA N o X)",
]

TEXTO_TIPO_PUESTO = [
    "Esta identificación difiere del tipo de HW a instalar, siendo así:",
    "     S - Puesto Sobremesa o Nettop",
    "     P - Puesto Portátil",
]


def agregar_slide_nomenclatura(prs, ruta_logo_color: str = None):
    slide = nueva_slide_en_blanco(prs)

    _caja_texto(slide, Inches(0.4), Inches(0.15), Inches(10), Inches(0.55),
               "Nomenclatura de equipos Colombia", tamano=24, color=ROJO_SANTANDER)

    _caja_texto(slide, Inches(0.4), Inches(0.75), Inches(12), Inches(0.5),
               TEXTO_INTRO_NOMENCLATURA, tamano=11, negrita=True, color=NEGRO)

    _caja_texto(slide, Inches(0.4), Inches(1.35), Inches(4), Inches(0.3),
               "EMPRESA (2-7 posiciones):", tamano=11, negrita=True, color=NEGRO)

    filas = len(TABLA_NOMENCLATURA)
    tabla_shape = slide.shapes.add_table(filas, 4, Inches(0.4), Inches(1.7), Inches(9), Inches(1.7))
    tabla = tabla_shape.table
    anchos = [Inches(3.6), Inches(1.3), Inches(2.1), Inches(2)]
    for i, ancho in enumerate(anchos):
        tabla.columns[i].width = ancho

    for fila_idx, fila in enumerate(TABLA_NOMENCLATURA):
        for col_idx, valor in enumerate(fila):
            celda = tabla.cell(fila_idx, col_idx)
            if fila_idx == 0:
                celda.fill.solid()
                celda.fill.fore_color.rgb = NAVY_BARRA
                _formatear_celda_color(celda, valor, tamano=11, negrita=True, color=BLANCO)
            else:
                celda.fill.solid()
                celda.fill.fore_color.rgb = BLANCO
                _formatear_celda_color(celda, valor, tamano=10.5, negrita=False, color=NEGRO,
                                       alineacion=PP_ALIGN.LEFT if col_idx == 0 else PP_ALIGN.CENTER)

    y_nota = Inches(3.55)
    _rect(slide, Inches(0.4), y_nota, Inches(9), Inches(0.45), color_relleno=GRIS_CLARO)
    caja_nota = slide.shapes.add_textbox(Inches(0.55), y_nota, Inches(8.7), Inches(0.45))
    tf = caja_nota.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = "Nota: "
    r1.font.bold = True
    r1.font.color.rgb = ROJO_SANTANDER
    r1.font.size = Pt(10.5)
    r1.font.name = FUENTE
    r2 = p.add_run()
    r2.text = NOTA_NOMENCLATURA.replace("Nota: ", "")
    r2.font.size = Pt(10.5)
    r2.font.name = FUENTE
    r2.font.color.rgb = NEGRO

    y = Inches(4.2)
    _caja_texto(slide, Inches(0.4), y, Inches(6), Inches(0.3),
               "MATRICULA DEL EMPLEADO O USUARIO:", tamano=11, negrita=True, color=NEGRO)
    y += Inches(0.35)
    _agregar_bullets(slide, Inches(0.6), y, Inches(11.5), Inches(1.1), TEXTO_MATRICULA)

    y = Inches(5.6)
    _caja_texto(slide, Inches(0.4), y, Inches(8), Inches(0.3),
               "TIPO DE PUESTO Y HARDWARE (1 posición al final del nombre de equipo):",
               tamano=11, negrita=True, color=NEGRO)
    y += Inches(0.35)
    _agregar_bullets(slide, Inches(0.6), y, Inches(11.5), Inches(1), TEXTO_TIPO_PUESTO)

    if ruta_logo_color:
        try:
            slide.shapes.add_picture(ruta_logo_color, Inches(0.4), Inches(6.9), height=Inches(0.4))
        except Exception:
            pass

    return slide


def _formatear_celda_color(celda, texto, tamano=10, negrita=False, color=NEGRO, alineacion=PP_ALIGN.CENTER):
    celda.margin_left = Pt(3)
    celda.margin_right = Pt(3)
    celda.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf = celda.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alineacion
    run = p.add_run()
    run.text = texto
    run.font.size = Pt(tamano)
    run.font.bold = negrita
    run.font.name = FUENTE
    run.font.color.rgb = color


def _agregar_bullets(slide, x, y, w, h, lineas: list[str]):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, linea in enumerate(lineas):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        prefijo = "• " if not linea.startswith("     ") else "   - "
        run.text = f"{prefijo}{linea.strip()}"
        run.font.size = Pt(10.5)
        run.font.name = FUENTE
        run.font.color.rgb = NEGRO